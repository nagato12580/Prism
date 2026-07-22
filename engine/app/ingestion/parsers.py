# engine/app/ingestion/parsers.py
from dataclasses import dataclass, field
from pathlib import Path
from typing import Protocol


class UnsupportedDocument(ValueError):
    pass


class ParserError(ValueError):
    """Typed error for parse failures (empty text, corruption, etc.)."""

    def __init__(self, parser_id: str, reason: str):
        self.parser_id = parser_id
        self.reason = reason
        super().__init__(f"{parser_id}: {reason}")


class DocumentParser(Protocol):
    parser_id: str
    extensions: frozenset[str]

    def parse(self, path: Path, config: dict) -> "ParsedDocument": ...


@dataclass(frozen=True)
class ParsedDocument:
    markdown: str
    parser_id: str
    page_count: int | None = None
    assets: tuple[str, ...] = ()
    metadata: dict = field(default_factory=dict)


class ParserRegistry:
    def __init__(self, parsers: list[DocumentParser]):
        self._parsers = parsers

    def parse(self, path: Path, media_type: str, config: dict) -> ParsedDocument:
        suffix = path.suffix.lower()
        parser = next((candidate for candidate in self._parsers if suffix in candidate.extensions), None)
        if parser is None:
            raise UnsupportedDocument(f"Unsupported file extension: {suffix}")
        return parser.parse(path, config)

    def capabilities(self) -> list[dict]:
        return [{"parser_id": p.parser_id, "extensions": sorted(p.extensions)} for p in self._parsers]


class MarkdownParser:
    parser_id = "markdown"
    extensions = frozenset({".md", ".markdown"})

    def parse(self, path: Path, config: dict) -> ParsedDocument:
        content = path.read_text(encoding="utf-8")
        if not content.strip():
            raise ParserError(self.parser_id, "empty document")
        return ParsedDocument(markdown=content, parser_id=self.parser_id)


class TextParser:
    parser_id = "text"
    extensions = frozenset({".txt"})

    def parse(self, path: Path, config: dict) -> ParsedDocument:
        content = path.read_text(encoding="utf-8")
        if not content.strip():
            raise ParserError(self.parser_id, "empty document")
        return ParsedDocument(markdown=content, parser_id=self.parser_id)


class PdfParser:
    parser_id = "pdf"
    extensions = frozenset({".pdf"})

    def parse(self, path: Path, config: dict) -> ParsedDocument:
        try:
            from backend.app.utils.file_parser import _extract_pdf, count_pages
        except ImportError:
            raise ParserError(self.parser_id, "PDF parser backend unavailable")
        text = _extract_pdf(str(path))
        if not text.strip():
            raise ParserError(self.parser_id, "empty or unreadable PDF")
        pages = count_pages(str(path)) if hasattr(count_pages, "__call__") else None
        return ParsedDocument(markdown=text, parser_id=self.parser_id, page_count=pages)


class DocxParser:
    parser_id = "docx"
    extensions = frozenset({".docx"})

    def parse(self, path: Path, config: dict) -> ParsedDocument:
        try:
            from backend.app.utils.file_parser import _extract_docx
        except ImportError:
            raise ParserError(self.parser_id, "DOCX parser backend unavailable")
        text = _extract_docx(str(path))
        if not text.strip():
            raise ParserError(self.parser_id, "empty or unreadable DOCX")
        return ParsedDocument(markdown=text, parser_id=self.parser_id)


class XlsxParser:
    parser_id = "xlsx"
    extensions = frozenset({".xlsx"})

    def parse(self, path: Path, config: dict) -> ParsedDocument:
        try:
            from backend.app.utils.file_parser import _extract_xlsx
        except ImportError:
            raise ParserError(self.parser_id, "XLSX parser backend unavailable")
        text = _extract_xlsx(str(path))
        if not text.strip():
            raise ParserError(self.parser_id, "empty or unreadable XLSX")
        return ParsedDocument(markdown=text, parser_id=self.parser_id)


class PptxParser:
    parser_id = "pptx"
    extensions = frozenset({".pptx"})

    def parse(self, path: Path, config: dict) -> ParsedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise ParserError(self.parser_id, "PPTX parser requires python-pptx")
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        text = "\n".join(texts)
        if not text.strip():
            raise ParserError(self.parser_id, "empty or unreadable PPTX")
        return ParsedDocument(markdown=text, parser_id=self.parser_id)


def build_default_registry() -> ParserRegistry:
    return ParserRegistry([
        MarkdownParser(),
        TextParser(),
        PdfParser(),
        DocxParser(),
        XlsxParser(),
        PptxParser(),
    ])
