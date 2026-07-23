# engine/tests/test_parser_fixtures.py
"""Tests that verify real PDF/DOCX/XLSX/PPTX parsing through the registry."""
import pytest
from pathlib import Path


# Minimal valid PDF — 14 bytes: header + cross-ref table + trailer + %%EOF
# This is a valid empty PDF that PyMuPDF/fitz can open.
_MINIMAL_PDF = (
    b"%PDF-1.4\n"
    b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
    b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
    b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]"
    b"/Contents 4 0 R/Resources<<>>>>endobj\n"
    b"4 0 obj<</Length 5 0 R>>stream\nBT /F1 12 Tf 100 700 Td (Hello PDF world) Tj ET\nendstream\nendobj\n"
    b"5 0 obj 44\nendobj\n"
    b"xref\n0 6\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n \n"
    b"0000000115 00000 n \n0000000205 00000 n \n0000000299 00000 n \n"
    b"trailer<</Size 6/Root 1 0 R>>\nstartxref\n329\n%%EOF"
)


def _generate_pdf(path: Path):
    path.write_bytes(_MINIMAL_PDF)


def _generate_docx(path: Path):
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx not installed")
    doc = Document()
    doc.add_paragraph("Hello DOCX world")
    doc.add_paragraph("Second paragraph")
    doc.save(str(path))


def _generate_xlsx(path: Path):
    try:
        from openpyxl import Workbook
    except ImportError:
        pytest.skip("openpyxl not installed")
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Hello"
    ws["B1"] = "XLSX"
    ws["A2"] = "Row2Col1"
    ws["B2"] = "Row2Col2"
    wb.save(str(path))


def _generate_pptx(path: Path):
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello PPTX world"
    prs.save(str(path))


def test_registry_parses_real_pdf(tmp_path):
    _generate_pdf(tmp_path / "test.pdf")
    from engine.app.ingestion.parsers import build_default_registry
    registry = build_default_registry()
    result = registry.parse(tmp_path / "test.pdf", media_type="document", config={})
    assert result.parser_id == "pdf"
    assert "Hello PDF world" in result.markdown


def test_registry_parses_real_docx(tmp_path):
    _generate_docx(tmp_path / "test.docx")
    from engine.app.ingestion.parsers import build_default_registry
    registry = build_default_registry()
    result = registry.parse(tmp_path / "test.docx", media_type="document", config={})
    assert result.parser_id == "docx"
    assert "Hello DOCX world" in result.markdown


def test_registry_parses_real_xlsx(tmp_path):
    _generate_xlsx(tmp_path / "test.xlsx")
    from engine.app.ingestion.parsers import build_default_registry
    registry = build_default_registry()
    result = registry.parse(tmp_path / "test.xlsx", media_type="spreadsheet", config={})
    assert result.parser_id == "xlsx"
    assert "Hello" in result.markdown
    assert "XLSX" in result.markdown


def test_registry_parses_real_pptx(tmp_path):
    _generate_pptx(tmp_path / "test.pptx")
    from engine.app.ingestion.parsers import build_default_registry
    registry = build_default_registry()
    result = registry.parse(tmp_path / "test.pptx", media_type="document", config={})
    assert result.parser_id == "pptx"
    assert "Hello PPTX world" in result.markdown


def test_parser_error_wraps_low_level_exception(tmp_path):
    """A corrupted file should produce a ParserError, not crash."""
    path = tmp_path / "corrupt.pdf"
    path.write_bytes(b"not a pdf")
    from engine.app.ingestion.parsers import ParserError, build_default_registry
    registry = build_default_registry()
    with pytest.raises(ParserError):
        registry.parse(path, media_type="document", config={})


def test_media_type_and_config_passed_to_parsers(tmp_path):
    _generate_pdf(tmp_path / "test.pdf")
    from engine.app.ingestion.parsers import build_default_registry
    registry = build_default_registry()
    result = registry.parse(
        tmp_path / "test.pdf", media_type="report", config={"page_count": 5}
    )
    assert result.parser_id == "pdf"
    assert result.page_count == 5
